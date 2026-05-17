"""Generate the OSS Preflight submission PowerPoint deck.

Content sourced from docs/submission-readiness.md (P9, 2026-05-17),
README.md, and docs/submission-deck.html. Brand palette mirrors
docs/submission-deck.html.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Brand palette
BG_DARK = RGBColor(0x0D, 0x13, 0x21)
BG_MID = RGBColor(0x1C, 0x2B, 0x36)
GOLD = RGBColor(0xF8, 0xB5, 0x52)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF2, 0xE9)
MUTED = RGBColor(0xAE, 0xBD, 0xB9)
PANEL = RGBColor(0x22, 0x33, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    # accent corner bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    bar.shadow.inherit = False
    return s


def txt(slide, left, top, width, height, text, size, color, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI", anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def kicker(slide, text):
    txt(slide, Inches(0.7), Inches(0.42), Inches(8), Inches(0.4),
        text.upper(), 15, GOLD, bold=True)


def footer(slide, section):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "OSS Preflight"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED
    r.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = "          •          " + section + "          •          IBM Bob Hackathon"
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED
    r2.font.name = "Segoe UI"


def panel(slide, left, top, width, height, title, body):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = RGBColor(0x3A, 0x4B, 0x52)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = GOLD
    r.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    p2.line_spacing = 1.15
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(13)
    r2.font.color.rgb = CREAM
    r2.font.name = "Segoe UI"


def metric(slide, left, top, width, big, small):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.fill.background()
    box.shadow.inherit = False
    # teal top rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()
    rule.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = big
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = small
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    r2.font.name = "Segoe UI"


def bullets(slide, left, top, width, height, items, size=17):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.2
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = TEAL
        rb.font.name = "Segoe UI"
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.color.rgb = CREAM
        r.font.name = "Segoe UI"
    return tb


# ---------- Slide 1: Title ----------
s = add_slide()
kicker(s, "IBM Bob Hackathon  —  Submission Deck")
txt(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.6),
    "OSS Preflight", 60, WHITE, bold=True)
txt(s, Inches(0.7), Inches(3.05), Inches(11.5), Inches(1.2),
    "Turn an app idea into a ranked, evidence-backed open source starting "
    "point — with a generated starter and a passing smoke test.",
    22, CREAM, spacing=1.25)
mw = Inches(2.9)
gap = Inches(0.18)
x0 = Inches(0.7)
for i, (b, sm) in enumerate([
    ("3", "ranked recommendations"),
    ("13/13", "production checks, 0 blockers"),
    ("181", "tests green (unmocked CLI)"),
    ("npm + PyPI", "real registry evidence"),
]):
    metric(s, x0 + i * (mw + gap), Inches(4.8), mw, b, sm)
footer(s, "Title")

# ---------- Slide 2: Problem ----------
s = add_slide()
kicker(s, "Problem")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.4),
    "Developers lose time choosing a stack before they write useful code.",
    34, WHITE, bold=True, spacing=1.05)
bullets(s, Inches(0.7), Inches(2.9), Inches(11.8), Inches(3.5), [
    "Package choice is driven by memory, popularity, or search ranking — not evidence.",
    "Measured facts, inferred tradeoffs, and stale assumptions get mixed together.",
    "Even a good choice still leaves starter setup and verification work undone.",
], size=19)
footer(s, "Problem")

# ---------- Slide 3: Solution ----------
s = add_slide()
kicker(s, "Solution")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "One honest pipeline: idea → evidence → scaffold → smoke test.",
    32, WHITE, bold=True)
pw = Inches(3.78)
pg = Inches(0.2)
py = Inches(2.5)
ph = Inches(3.6)
panel(s, Inches(0.7), py, pw, ph, "Recommend",
      "Parse the idea, run search-first discovery across npm / PyPI / GitHub, "
      "collect public facts, and rank candidates deterministically.")
panel(s, Inches(0.7) + pw + pg, py, pw, ph, "Explain",
      "Every package gets an Evidence Passport that separates measured "
      "registry facts from AI-generated interpretation.")
panel(s, Inches(0.7) + 2 * (pw + pg), py, pw, ph, "Generate",
      "Scaffold a runnable starter with mocked messages and a smoke test, "
      "or emit an adoption pack when no template fits.")
footer(s, "Solution")

# ---------- Slide 4: Evidence Passport ----------
s = add_slide()
kicker(s, "Differentiator")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "The Evidence Passport: measured fact vs. assessed interpretation.",
    30, WHITE, bold=True)
panel(s, Inches(0.7), Inches(2.5), Inches(5.85), Inches(3.7), "Measured (left)",
      "Registry metadata, license, latest version, last commit, downloads, "
      "OpenSSF score, and the fetch timestamp. Every fact carries its source "
      "and collectedAt; missing data is explicit not-available.")
panel(s, Inches(6.78), Inches(2.5), Inches(5.85), Inches(3.7), "Assessed (right)",
      "Goal fit, tradeoffs, and warnings. Clearly labeled as OSS Preflight's "
      "interpretation. The measured/assessed split is enforced in the data "
      "model — not just the UI.")
footer(s, "Evidence Passport")

# ---------- Slide 5: How it works ----------
s = add_slide()
kicker(s, "Architecture")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "Search-first discovery with an honest workflow trace.",
    30, WHITE, bold=True)
bullets(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(4.0), [
    "Discovery calls live npm/PyPI/GitHub search, then blends labeled catalog "
    "fallback so broad registry noise can't crowd out known-fit packages.",
    "Every run writes workflow.json recording the real discovery method "
    "(registry-search / search-with-catalog-fallback / manifest) and candidates.",
    "Three real user paths: idea → recommend → scaffold/adoption → run; "
    "audit an existing npm or Python project; same CLI via Web bridge or Bob skill.",
    "BYOK intent parsing (Anthropic / OpenAI-compatible / Gemini) with a "
    "deterministic keyword fallback — the demo runs on a fresh clone, no keys.",
], size=16)
footer(s, "Architecture")

# ---------- Slide 6: Demo ----------
s = add_slide()
kicker(s, "Demo")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "The verified local demo path is ready.",
    32, WHITE, bold=True)
bullets(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(3.8), [
    "Input:  “Discord bot that summarizes channel activity”",
    "Result:  three recommendations, with discord.js ranked #1.",
    "Scaffold:  generated starter installs deps and passes its smoke test "
    "(mocked Discord messages, no credentials required).",
    "Build Proof tab surfaces Bob modes, skills, and session exports.",
    "Fallback:  README documents the full local path if no public URL is live.",
], size=17)
footer(s, "Demo")

# ---------- Slide 7: Validation & quality ----------
s = add_slide()
kicker(s, "Evidence & Quality")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "Every claim is backed by an executed command.",
    30, WHITE, bold=True)
txt(s, Inches(0.7), Inches(2.05), Inches(11.8), Inches(0.7),
    "pnpm validate:production runs the real product — build, full test "
    "suite, and live CLI — and fails loudly on regression.",
    15, CREAM, spacing=1.2)
data = [
    ("Agentic workflow maturity", "8/10", "was ~6"),
    ("Discord CLI demo reliability", "9/10", "was ~8"),
    ("Full product demo reliability", "8/10", "was ~6"),
    ("Usefulness on arbitrary projects", "8/10", "was ~4"),
    ("Hackathon evidence clarity", "8/10", "was ~5"),
]
ry = Inches(3.05)
rh = Inches(0.62)
for i, (label, score, base) in enumerate(data):
    y = ry + i * (rh + Inches(0.08))
    row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(11.93), rh)
    row.fill.solid()
    row.fill.fore_color.rgb = PANEL if i % 2 == 0 else BG_MID
    row.line.fill.background()
    row.shadow.inherit = False
    txt(s, Inches(0.95), y + Inches(0.13), Inches(7.5), Inches(0.4),
        label, 15, CREAM)
    txt(s, Inches(8.5), y + Inches(0.13), Inches(1.6), Inches(0.4),
        score, 16, GOLD, bold=True)
    txt(s, Inches(10.2), y + Inches(0.15), Inches(2.2), Inches(0.4),
        base, 12, MUTED)
footer(s, "Evidence & Quality")

# ---------- Slide 8: Bob evidence ----------
s = add_slide()
kicker(s, "Bob Evidence")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.1),
    "Bob is part of the SDLC, not just the final polish.",
    30, WHITE, bold=True)
pw = Inches(3.78)
pg = Inches(0.2)
py = Inches(2.5)
ph = Inches(3.6)
panel(s, Inches(0.7), py, pw, ph, "Modes",
      "Plan, Code, Orchestrator, reviewer, and scoped scaffolder modes drive "
      "the build with checkpoints.")
panel(s, Inches(0.7) + pw + pg, py, pw, ph, "Skills",
      "Code review, evidence discipline, test runner, doc writer, Plan "
      "Council, and the OSS Preflight advisor skill.")
panel(s, Inches(0.7) + 2 * (pw + pg), py, pw, ph, "Exports",
      "A single build ledger plus per-session folders committed under "
      "bob_sessions/ as machine-checkable evidence.")
footer(s, "Bob Evidence")

# ---------- Slide 9: Team + Readiness ----------
s = add_slide()
kicker(s, "Team & Readiness")
txt(s, Inches(0.7), Inches(1.05), Inches(11.8), Inches(1.3),
    "Solo build, Bob-assisted delivery, honest remaining actions.",
    30, WHITE, bold=True)
bullets(s, Inches(0.7), Inches(2.7), Inches(11.9), Inches(2.6), [
    "Team:  Dreamverse Holdings solo builder, with IBM Bob as the visible engineering partner.",
    "Green:  pnpm test (181), pnpm build, pnpm validate:production (13/13, 0 blockers).",
    "Known gaps:  browser e2e is still a mocked UI smoke; only Discord has a code template.",
    "Still external:  demo video, missing Bob export screenshots, final approved tag.",
], size=17)
txt(s, Inches(0.7), Inches(5.7), Inches(11.8), Inches(0.6),
    "Phase P9 — Production Readiness  •  Reviewed 2026-05-17",
    14, GOLD, bold=True)
footer(s, "Team & Readiness")

out = r"docs\oss-preflight-submission-deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
